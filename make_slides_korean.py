"""
10슬라이드 한국어 버전 — 코리아 디스카운트 발표자료
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xD9, 0xD9, 0xD9)
MGRAY = RGBColor(0x59, 0x59, 0x59)
DGRAY = RGBColor(0x1A, 0x1A, 0x1A)

TITLE_FONT = "맑은 고딕"
BODY_FONT  = "맑은 고딕"
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

FIGURES = "/Users/dandan/Desktop/Projects/kor-discount/output/figures/png"


# ── 헬퍼 함수 ──────────────────────────────────────────────────────────────

def slide():
    return prs.slides.add_slide(BLANK)

def txb(sl, l, t, w, h):
    s = sl.shapes.add_textbox(l, t, w, h)
    s.text_frame.word_wrap = True
    return s.text_frame

def title_bar(sl, text, sub=None):
    r = sl.shapes.add_shape(1, 0, 0, W, Inches(0.07))
    r.fill.solid(); r.fill.fore_color.rgb = BLACK; r.line.fill.background()

    tf = txb(sl, Inches(0.55), Inches(0.15), Inches(12.2), Inches(0.75))
    p = tf.paragraphs[0]; run = p.add_run(); run.text = text
    run.font.name = TITLE_FONT; run.font.bold = True
    run.font.size = Pt(28); run.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.LEFT

    if sub:
        tf2 = txb(sl, Inches(0.55), Inches(0.85), Inches(12.2), Inches(0.38))
        p2 = tf2.paragraphs[0]; r2 = p2.add_run(); r2.text = sub
        r2.font.name = BODY_FONT; r2.font.size = Pt(14)
        r2.font.color.rgb = MGRAY

    br = sl.shapes.add_shape(1, 0, Inches(7.38), W, Inches(0.05))
    br.fill.solid(); br.fill.fore_color.rgb = BLACK; br.line.fill.background()

def para(tf, text, size=15, bold=False, italic=False,
         color=BLACK, space=0, indent=0):
    p = tf.add_paragraph()
    p.space_before = Pt(space)
    if indent: p.level = indent
    r = p.add_run(); r.text = text
    r.font.name = BODY_FONT; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return p

def bullet(tf, text, size=14, space=3, color=BLACK):
    para(tf, "•  " + text, size=size, space=space, indent=1, color=color)

def hline(sl, top):
    ln = sl.shapes.add_shape(1, Inches(0.55), top, Inches(12.2), Pt(1))
    ln.fill.solid(); ln.fill.fore_color.rgb = LGRAY; ln.line.fill.background()

def img(sl, name, l, t, w, h):
    sl.shapes.add_picture(f"{FIGURES}/{name}", l, t, w, h)

def body(sl, top=Inches(1.35)):
    return txb(sl, Inches(0.55), top, Inches(12.2), H - top - Inches(0.2))

def table_row(sl, cols, lefts, widths, top, height, is_header=False, shade=False):
    for text, cl, cw in zip(cols, lefts, widths):
        bg = sl.shapes.add_shape(1, cl, top, cw, height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BLACK if is_header else (LGRAY if shade else WHITE)
        bg.line.fill.background()
        tf = txb(sl, cl + Pt(5), top + Pt(4), cw - Pt(10), height - Pt(5))
        p = tf.paragraphs[0]; r = p.add_run(); r.text = text
        r.font.name = BODY_FONT; r.font.size = Pt(13)
        r.font.bold = is_header
        r.font.color.rgb = WHITE if is_header else BLACK
        p.alignment = PP_ALIGN.LEFT if lefts.index(cl) == 0 else PP_ALIGN.CENTER


# ════════════════════════════════════════════════════════════════════════════
#  슬라이드 1 — 표지
# ════════════════════════════════════════════════════════════════════════════
sl = slide()
bar = sl.shapes.add_shape(1, 0, 0, W, Inches(3.0))
bar.fill.solid(); bar.fill.fore_color.rgb = BLACK; bar.line.fill.background()

tf = txb(sl, Inches(0.7), Inches(0.4), Inches(11.9), Inches(2.3))
p = tf.paragraphs[0]; r = p.add_run()
r.text = "기업지배구조 개혁과 코리아 디스카운트"
r.font.name = TITLE_FONT; r.font.size = Pt(34); r.font.bold = True
r.font.color.rgb = WHITE

p2 = tf.add_paragraph(); p2.space_before = Pt(6)
r2 = p2.add_run(); r2.text = "일본 개혁 경험으로부터의 교훈"
r2.font.name = TITLE_FONT; r2.font.size = Pt(21); r2.font.color.rgb = LGRAY

meta = txb(sl, Inches(0.7), Inches(3.2), Inches(11.9), Inches(4.0))
for lbl, val in [
    ("분석 대상",  "월별 지수 PBR: KOSPI · TOPIX · S&P 500 · MSCI EM  |  2004년 1월 – 2024년 12월  |  관측치 1,008개"),
    ("핵심 결과",  "KOSPI PBR은 TOPIX 대비 평균 –0.177배 (t = –3.23), MSCI EM 대비 –0.601배 (t = –10.30) 낮음"),
    ("접근 방법", "일본의 3단계 지배구조 개혁을 정책 벤치마크로 활용; 이벤트 윈도우 · 패널 OLS · 합성통제법"),
    ("주의사항",   "국가 클러스터 N = 4; 포화 이벤트 윈도우 설계 — 기술적 분석이며 인과관계 식별 아님"),
]:
    p = meta.add_paragraph(); p.space_before = Pt(8)
    rl = p.add_run(); rl.text = lbl + ":  "
    rl.font.name = BODY_FONT; rl.font.size = Pt(14); rl.font.bold = True; rl.font.color.rgb = BLACK
    rv = p.add_run(); rv.text = val
    rv.font.name = BODY_FONT; rv.font.size = Pt(14); rv.font.color.rgb = BLACK

br = sl.shapes.add_shape(1, 0, Inches(7.38), W, Inches(0.05))
br.fill.solid(); br.fill.fore_color.rgb = BLACK; br.line.fill.background()


# ════════════════════════════════════════════════════════════════════════════
#  슬라이드 2 — 코리아 디스카운트
# ════════════════════════════════════════════════════════════════════════════
sl = slide()
title_bar(sl, "코리아 디스카운트",
          sub="KOSPI 주가는 20년 이상 비교 가능한 글로벌 지수 대비 지속적으로 저평가")

tf = body(sl, top=Inches(1.35))
para(tf, "할인 규모 — 2004–2024년 평균 PBR", size=16, bold=True)

COLS_L = [Inches(0.55), Inches(3.5), Inches(6.3), Inches(8.7), Inches(10.8)]
COLS_W = [Inches(2.9),  Inches(2.7), Inches(2.35),Inches(2.05),Inches(2.2)]
rows = [
    ("비교 지수", "평균 격차", "t-통계량", "95% 신뢰구간", "KOSPI 하위 기간"),
    ("vs. TOPIX",   "–0.177배", "–3.23",  "[–0.284배, –0.069배]", ""),
    ("vs. MSCI EM", "–0.601배", "–10.30", "[–0.716배, –0.486배]", ""),
    ("vs. S&P 500", "–1.90배",  "—",      "—",                    ""),
]

r_top = Inches(1.95); r_h = Inches(0.47)
for ri, row in enumerate(rows):
    table_row(sl, row, COLS_L, COLS_W, r_top + ri * r_h, r_h,
              is_header=(ri == 0), shade=(ri % 2 == 0))

note_tf = txb(sl, COLS_L[0], r_top + r_h + Pt(4), Inches(12.2), Pt(22))
note_p = note_tf.paragraphs[0]; note_r = note_p.add_run()
note_r.text = "KOSPI 추이: 1.36배 (2004–2013)  →  1.02배 (2014–2022)  →  0.93배 (2023–2024)  |  TOPIX는 2023년 TSE 개혁 후 1.36배로 반등"
note_r.font.name = BODY_FONT; note_r.font.size = Pt(12)
note_r.font.italic = True; note_r.font.color.rgb = MGRAY

hline(sl, Inches(3.5))

para(tf, "\n세 가지 구조적 원인", size=16, bold=True, space=4)
for ch in [
    "재벌 불투명성 및 대리인 비용 — 지배권-현금흐름권 비율 ≈ 2:1; 터널링; 주가 억제 유인(상속세 최대 60%)",
    "소수주주 권리 보호 미흡 — 민법 전통, 취약한 대표소송 권한, 스튜어드십 코드 실효성 부족",
    "지정학적 리스크 프리미엄 — 북한 핵실험 6회; 개별 이벤트 반응이 아닌 구조적 수준 할인으로 반영",
]:
    bullet(tf, ch, size=14)


# ════════════════════════════════════════════════════════════════════════════
#  슬라이드 3 — 그림 1: PBR 비교
# ════════════════════════════════════════════════════════════════════════════
sl = slide()
title_bar(sl, "그림 1: 지수별 PBR 추이 (2004–2024)",
          sub="점선은 일본의 3차례 개혁 시점; 2023년 이후 코리아 디스카운트 확대")

img(sl, "figure1_pb_comparison.png",
    Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.9))


# ════════════════════════════════════════════════════════════════════════════
#  슬라이드 4 — 일본의 지배구조 개혁 프로그램
# ════════════════════════════════════════════════════════════════════════════
sl = slide()
title_bar(sl, "일본의 지배구조 개혁 프로그램",
          sub="3단계 정책 개입 — 정책 벤치마크로 활용 (자연실험으로 간주하지 않음)")

reforms = [
    ("2014년 2월", "스튜어드십 코드",
     "기관투자자에게 참여 정책 수립 및 의결권 행사 내역 공시 의무화. "
     "Comply-or-explain 방식으로 높은 이행률 달성. "
     "케이레츠 교차보유 해소 촉진 (Miyajima 2023)."),
    ("2015년 6월", "기업지배구조 코드",
     "모든 TSE 상장사에 지배구조 원칙 준수 또는 미준수 이유 공시 의무화. "
     "사외이사 2인 이상 선임 요구. 2017년까지 JPX-Nikkei 400 구성사의 90% 이상 이행."),
    ("2023년 3월", "TSE PBR 개혁",
     "PBR 1.0배 미만 기업에 자본 효율 개선 계획 공시 요구. "
     "TOPIX PBR ≈1.0배 → ≈1.36배 상승. KOSPI는 ≈0.93배 유지 → 한·일 격차 확대. "
     "한국의 기업가치 제고 프로그램(2024)은 이 개혁을 모델로 하나 자율 참여 방식."),
]

col_l = [Inches(0.55), Inches(1.9), Inches(4.1)]
col_w = [Inches(1.3),  Inches(2.15),Inches(9.0)]
r_top = Inches(1.35); r_h = Inches(1.65)

for ri, (date, name, desc) in enumerate(reforms):
    for ci, (cl, cw, fg) in enumerate(zip(col_l, col_w,
                                          [BLACK, DGRAY, LGRAY if ri % 2 == 0 else WHITE])):
        bg = sl.shapes.add_shape(1, cl, r_top + ri * r_h, cw, r_h)
        bg.fill.solid(); bg.fill.fore_color.rgb = fg; bg.line.fill.background()

    for cl, cw, txt, fc, fsz, fbd in [
        (col_l[0], col_w[0], date, WHITE, 15, True),
        (col_l[1], col_w[1], name, WHITE, 13, True),
        (col_l[2], col_w[2], desc, BLACK, 12, False),
    ]:
        tf2 = txb(sl, cl + Pt(5), r_top + ri * r_h + Pt(7), cw - Pt(10), r_h - Pt(10))
        p2 = tf2.paragraphs[0]; r2 = p2.add_run(); r2.text = txt
        r2.font.name = BODY_FONT; r2.font.size = Pt(fsz)
        r2.font.bold = fbd; r2.font.color.rgb = fc

note = txb(sl, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.55))
note.paragraphs[0].add_run().text = (
    "한국 비교: 스튜어드십 코드 2016년 12월 도입(이행 강도 약함), 기업가치 제고 프로그램 2024년 도입(자율 참여) — "
    "완전한 미처리 비교군이 아닌 부분 비교군. "
    "모든 개혁 시점은 데이터 변환 이전 config.py에 고정 (look-ahead bias 방지).")
note.paragraphs[0].runs[0].font.size = Pt(11)
note.paragraphs[0].runs[0].font.italic = True
note.paragraphs[0].runs[0].font.color.rgb = MGRAY


# ════════════════════════════════════════════════════════════════════════════
#  슬라이드 5 — 데이터 및 실증 전략
# ════════════════════════════════════════════════════════════════════════════
sl = slide()
title_bar(sl, "데이터 및 실증 분석 전략")

ltf = txb(sl, Inches(0.55), Inches(1.35), Inches(5.8), Inches(5.85))
para(ltf, "데이터", size=17, bold=True)
for d in [
    "Bloomberg PX_TO_BOOK_RATIO: 월별 지수 PBR",
    "KOSPI · TOPIX · S&P 500 · MSCI EM",
    "2004년 1월 – 2024년 12월 (국가-월 관측치 1,008개)",
    "PBR 선택 이유: 경기 변동 민감도 낮음; 불황기 음수값 없음; TSE 개혁의 직접 대상",
    "지정학적 리스크 지수(Caldara–Iacoviello) + 북한 위협 GPRNK 하위지수",
    "4개 시계열 모두 결측값 없음",
]:
    bullet(ltf, d, size=13)

vl = sl.shapes.add_shape(1, Inches(6.55), Inches(1.35), Pt(1), Inches(5.85))
vl.fill.solid(); vl.fill.fore_color.rgb = LGRAY; vl.line.fill.background()

rtf = txb(sl, Inches(6.75), Inches(1.35), Inches(6.1), Inches(5.85))
para(rtf, "세 가지 실증 분석 방법", size=17, bold=True)

for name, desc in [
    ("이벤트 윈도우 분석  [주요 분석]",
     "[–36, –1]월 구간에서 선형 사전 추세 추출; CAR = [–12, +24] 구간 KOSPI–TOPIX 스프레드의 누적 비정상 변화. "
     "포화 설계 → 표준오차 미정의. 결과는 기술적 분석에 한함."),
    ("패널 OLS  [보완 분석]",
     "이원 고정효과(국가 + 시간); 개혁 × 일본 교호항. "
     "와일드 부트스트랩 p값 (Rademacher 999회 추출, N = 4 클러스터). "
     "결과: p = 0.625 / 0.375 / 0.500 — 강건한 인과적 증거 없음."),
    ("합성통제법  [2023년 한정]",
     "역내 EM 비교군에 대한 디민(demeaned) PBR 기반 합성통제. "
     "MSCI EM 아시아에 전체 가중치 배분(w = 1.0); 처치 전 RMSPE = 0.1451. "
     "사후 격차는 확대되지 않고 점진적으로 축소."),
]:
    para(rtf, name, size=14, bold=True, space=8)
    para(rtf, desc, size=13, space=2, indent=1)


# ════════════════════════════════════════════════════════════════════════════
#  슬라이드 6 — 이벤트 윈도우 분석 결과 + 그림 2
# ════════════════════════════════════════════════════════════════════════════
sl = slide()
title_bar(sl, "결과: 이벤트 윈도우 분석",
          sub="누적 비정상 KOSPI–TOPIX PBR 스프레드; 양수 = 사전 추세 대비 할인 축소")

img(sl, "figure2_event_study.png",
    Inches(0.3), Inches(1.35), Inches(7.5), Inches(5.85))

rtf = txb(sl, Inches(8.0), Inches(1.35), Inches(5.0), Inches(5.85))
para(rtf, "τ = +24에서의 코호트별 CAR", size=15, bold=True)

for cohort, car, note in [
    ("2014 스튜어드십 코드",    "CAR ≈ +4.29",
     "스프레드가 사전 추세 대비 상승. 2015년 윈도우와 중복 — 독립적 실험 아님."),
    ("2015 기업지배구조 코드",  "CAR ≈ +7.83",
     "3개 코호트 중 가장 큰 누적 움직임. 2014년 윈도우와 중복."),
    ("2023 TSE PBR 개혁",       "CAR ≈ –6.48",
     "개혁 후 한·일 격차가 사전 추세 대비 확대. "
     "정책적으로 가장 시의성 있는 코호트."),
]:
    para(rtf, cohort + "  |  " + car, size=14, bold=True, space=10)
    para(rtf, note, size=12, space=1, indent=1, color=MGRAY)

para(rtf,
    "\n주의: 설계가 포화 상태이므로 신뢰구간 없음; "
    "2014·2015 코호트는 중복됨. 인과적 처치 효과가 아닌 "
    "기술적 스프레드 변화임.",
    size=12, italic=True, color=MGRAY, space=10)


# ════════════════════════════════════════════════════════════════════════════
#  슬라이드 7 — 패널 OLS + 합성통제법
# ════════════════════════════════════════════════════════════════════════════
sl = slide()
title_bar(sl, "결과: 패널 OLS 및 합성통제법")

ltf = txb(sl, Inches(0.55), Inches(1.35), Inches(6.0), Inches(5.85))
para(ltf, "패널 OLS — 이원 고정효과 + 개혁 × 일본", size=16, bold=True)
para(ltf, "와일드 부트스트랩 p값 (Rademacher 999회, N = 4 클러스터)", size=12,
     italic=True, color=MGRAY, space=2)

ols_rows = [
    ("교호 변수",                         "계수",   "p값"),
    ("스튜어드십 코드 × 일본",            "+0.09",  "0.625"),
    ("기업지배구조 코드 × 일본",          "–0.32",  "0.375"),
    ("TSE PBR 개혁 × 일본",               "–0.23",  "0.500"),
]
cl = [Inches(0.55), Inches(4.0), Inches(5.15)]
cw = [Inches(3.4),  Inches(1.1), Inches(1.35)]
rt = Inches(2.35); rh = Inches(0.48)

for ri, row in enumerate(ols_rows):
    table_row(sl, row, cl, cw, rt + ri * rh, rh, is_header=(ri == 0), shade=(ri % 2 == 0))

para(ltf, "\n추정치 3개 중 2개가 음수.\n와일드 부트스트랩에서 p < 0.05 달성 없음.\nN = 4 클러스터 → 통계적 검정력 매우 제한적.\n해석: 강건한 인과적 증거 없음.",
     size=13, space=6)

vl = sl.shapes.add_shape(1, Inches(6.7), Inches(1.35), Pt(1), Inches(5.85))
vl.fill.solid(); vl.fill.fore_color.rgb = LGRAY; vl.line.fill.background()

rtf = txb(sl, Inches(6.9), Inches(1.35), Inches(6.1), Inches(5.85))
para(rtf, "합성통제법 — 2023 TSE PBR 개혁", size=16, bold=True)
para(rtf, "비교군: MSCI EM 아시아, MSCI 홍콩, MSCI 대만, MSCI 중국, MSCI 인도, MSCI 인도네시아",
     size=12, italic=True, color=MGRAY, space=2)

for lbl, val in [
    ("비교군 가중치",     "MSCI EM 아시아 = 1.0  (단일 비교군 해)"),
    ("처치 전 RMSPE",     "0.1451  (임계값 0.15 이하)"),
    ("사후 격차",         "소폭 양수이나 점진적 축소 (월평균 변화 ≈ –0.0028)"),
    ("해석",              "일본의 PBR 상승은 합성 EM 아시아 비교군 추이와 대체로 일치. "
                          "2023년 이후 독립적 지배구조 프리미엄 지속되지 않음."),
]:
    p = rtf.add_paragraph(); p.space_before = Pt(9)
    r1 = p.add_run(); r1.text = lbl + ":  "
    r1.font.name = BODY_FONT; r1.font.size = Pt(14); r1.font.bold = True; r1.font.color.rgb = BLACK
    r2 = p.add_run(); r2.text = val
    r2.font.name = BODY_FONT; r2.font.size = Pt(14); r2.font.color.rgb = BLACK

para(rtf,
    "\n주의: TOPIX의 절대 PBR 상승은 엔화 약세 및 "
    "역내 EM 모멘텀을 반영할 수 있으며, "
    "순수한 지배구조 프리미엄으로 단정할 수 없음.",
    size=12, italic=True, color=MGRAY, space=8)


# ════════════════════════════════════════════════════════════════════════════
#  슬라이드 8 — 그림: 지정학적 리스크 + 합성통제 격차
# ════════════════════════════════════════════════════════════════════════════
sl = slide()
title_bar(sl, "예비 결과: 추가 그림",
          sub="좌: 지정학적 리스크 계수 유의하지 않음 (p = 0.40) — 구조적 수준 할인으로 반영   |   우: 합성통제 격차는 2023년 이후 점차 축소")

img(sl, "figure3_geo_risk.png",
    Inches(0.3), Inches(1.35), Inches(6.3), Inches(5.85))

img(sl, "figure_synth_gap.png",
    Inches(6.9), Inches(1.35), Inches(6.3), Inches(5.85))


# ════════════════════════════════════════════════════════════════════════════
#  슬라이드 9 — 반사실적 전망
# ════════════════════════════════════════════════════════════════════════════
sl = slide()
title_bar(sl, "예시적 반사실적 전망",
          sub="한국이 지배구조 개혁을 시행할 경우 KOSPI PBR 경로 — 민감도 분석이며 예측이 아님")

img(sl, "figure4_counterfactual_projection.png",
    Inches(0.3), Inches(1.35), Inches(7.5), Inches(5.85))

rtf = txb(sl, Inches(8.0), Inches(1.35), Inches(5.0), Inches(5.85))
para(rtf, "산출 방법", size=16, bold=True)
for d in [
    "기준: KOSPI PBR ≈ 0.93배 (2024년 12월)",
    "적용 변화율: 2023년 TSE 개혁 후 18개월간 일본 합성통제 격차의 월평균 변화량 (–0.0028/월)",
    "불확실성 구간: ±RMSPE = ±0.1451",
]:
    bullet(rtf, d, size=13)

para(rtf, "\n결과", size=16, bold=True, space=10)
para(rtf,
    "일본의 2023년 이후 상대 격차 변화가 음수이므로, "
    "기계적 전망은 코리아 디스카운트를 해소하지 못함. "
    "정책 성공 예측이 아닌 민감도 분석임.",
    size=13, space=2)

para(rtf, "\n정책 레버 (개혁 시행 시)", size=15, bold=True, space=10)
for rec in [
    "금융위원회: PBR 공시 의무화 (자율 참여 → comply-or-explain)",
    "한국거래소: 대형 상장사 사외이사 비율 50% 이상",
    "스튜어드십 코드: 의결권 행사 내역 의무 공시 + 국민연금 거버넌스 전담 조직 신설",
]:
    bullet(rtf, rec, size=13)


# ════════════════════════════════════════════════════════════════════════════
#  슬라이드 10 — 결론
# ════════════════════════════════════════════════════════════════════════════
sl = slide()
title_bar(sl, "결론")

tf = body(sl, top=Inches(1.35))

for heading, bullets in [
    ("주요 발견",
     [
         "KOSPI PBR은 TOPIX 대비 평균 –0.177배 (t = –3.23), MSCI EM 대비 –0.601배 (t = –10.30)",
         "2023년 이후 할인 확대: KOSPI ≈ 0.93배 vs. TOPIX ≈ 1.36배",
         "3대 구조적 원인: 재벌 불투명성 · 소수주주 보호 미흡 · 지정학적 리스크 프리미엄",
     ]),
    ("실증 결과",
     [
         "이벤트 윈도우 CAR: 개혁 시점 전후 스프레드 큰 변동 있으나 유효한 추론 불가 (포화 설계, 중복 코호트)",
         "패널 OLS p값: 0.625 / 0.375 / 0.500 — 일본 지배구조 재평가의 강건한 인과적 증거 없음",
         "합성통제법: MSCI EM 아시아 전체 가중(RMSPE 0.1451); 사후 격차 축소, 확대 아님",
         "지정학적 리스크 계수: –0.02 (p = 0.40) — 구조적 수준 할인으로 반영",
     ]),
    ("시사점",
     [
         "코리아 디스카운트는 지속적이고 통계적으로 유의미함",
         "일본의 경험은 유용한 정책 벤치마크를 제공하나, 한국이 일본식 개혁을 도입했을 때의 "
           "밸류에이션 개선 폭을 특정하지는 못함",
         "향후 연구: 재벌 지배구조 기업 패널 분석; 기업가치 제고 프로그램 이벤트 스터디 (2024–2026+)",
     ]),
]:
    para(tf, heading, size=16, bold=True, space=8)
    for b in bullets:
        bullet(tf, b, size=13)


# ════════════════════════════════════════════════════════════════════════════
#  저장
# ════════════════════════════════════════════════════════════════════════════
out = "/Users/dandan/Desktop/Projects/kor-discount/output/korea_discount_slides_korean.pptx"
prs.save(out)
print(f"저장 완료 → {out}")
