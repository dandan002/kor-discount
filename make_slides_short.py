"""
Condensed 10-slide version of the Korea Discount presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xD9, 0xD9, 0xD9)
MGRAY = RGBColor(0x59, 0x59, 0x59)
DGRAY = RGBColor(0x1A, 0x1A, 0x1A)

TITLE_FONT = "Calibri"
BODY_FONT  = "Calibri"
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

FIGURES = "/Users/dandan/Desktop/Projects/kor-discount/output/figures/png"


# ── helpers ──────────────────────────────────────────────────────────────────

def slide():
    return prs.slides.add_slide(BLANK)

def txb(sl, l, t, w, h):
    s = sl.shapes.add_textbox(l, t, w, h)
    s.text_frame.word_wrap = True
    return s.text_frame

def title_bar(sl, text, sub=None):
    r = sl.shapes.add_shape(1, 0, 0, W, Inches(0.07))
    r.fill.solid(); r.fill.fore_color.rgb = BLACK; r.line.fill.background()

    tf = txb(sl, Inches(0.55), Inches(0.15), Inches(12.2), Inches(0.75))
    p = tf.paragraphs[0]; run = p.add_run(); run.text = text
    run.font.name = TITLE_FONT; run.font.bold = True
    run.font.size = Pt(28); run.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.LEFT

    if sub:
        tf2 = txb(sl, Inches(0.55), Inches(0.85), Inches(12.2), Inches(0.38))
        p2 = tf2.paragraphs[0]; r2 = p2.add_run(); r2.text = sub
        r2.font.name = BODY_FONT; r2.font.size = Pt(14)
        r2.font.color.rgb = MGRAY

    br = sl.shapes.add_shape(1, 0, Inches(7.38), W, Inches(0.05))
    br.fill.solid(); br.fill.fore_color.rgb = BLACK; br.line.fill.background()

def para(tf, text, size=15, bold=False, italic=False,
         color=BLACK, space=0, indent=0):
    p = tf.add_paragraph()
    p.space_before = Pt(space)
    if indent: p.level = indent
    r = p.add_run(); r.text = text
    r.font.name = BODY_FONT; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return p

def bullet(tf, text, size=14, space=3, color=BLACK):
    para(tf, "•  " + text, size=size, space=space, indent=1, color=color)

def hline(sl, top):
    ln = sl.shapes.add_shape(1, Inches(0.55), top, Inches(12.2), Pt(1))
    ln.fill.solid(); ln.fill.fore_color.rgb = LGRAY; ln.line.fill.background()

def img(sl, name, l, t, w, h):
    sl.shapes.add_picture(f"{FIGURES}/{name}", l, t, w, h)

def body(sl, top=Inches(1.35)):
    return txb(sl, Inches(0.55), top, Inches(12.2), H - top - Inches(0.2))

def table_row(sl, cols, lefts, widths, top, height, is_header=False, shade=False):
    for text, cl, cw in zip(cols, lefts, widths):
        bg = sl.shapes.add_shape(1, cl, top, cw, height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BLACK if is_header else (LGRAY if shade else WHITE)
        bg.line.fill.background()
        tf = txb(sl, cl + Pt(5), top + Pt(4), cw - Pt(10), height - Pt(5))
        p = tf.paragraphs[0]; r = p.add_run(); r.text = text
        r.font.name = BODY_FONT; r.font.size = Pt(13 if not is_header else 13)
        r.font.bold = is_header
        r.font.color.rgb = WHITE if is_header else BLACK
        p.alignment = PP_ALIGN.LEFT if lefts.index(cl) == 0 else PP_ALIGN.CENTER


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
sl = slide()
bar = sl.shapes.add_shape(1, 0, 0, W, Inches(3.0))
bar.fill.solid(); bar.fill.fore_color.rgb = BLACK; bar.line.fill.background()

tf = txb(sl, Inches(0.7), Inches(0.4), Inches(11.9), Inches(2.3))
p = tf.paragraphs[0]; r = p.add_run()
r.text = "Corporate Governance Reform and the Korea Discount"
r.font.name = TITLE_FONT; r.font.size = Pt(34); r.font.bold = True
r.font.color.rgb = WHITE

p2 = tf.add_paragraph(); p2.space_before = Pt(6)
r2 = p2.add_run(); r2.text = "Lessons from Japan's Reform Experience"
r2.font.name = TITLE_FONT; r2.font.size = Pt(21); r2.font.color.rgb = LGRAY

meta = txb(sl, Inches(0.7), Inches(3.2), Inches(11.9), Inches(4.0))
for lbl, val in [
    ("Setting",  "Monthly index P/B ratios: KOSPI · TOPIX · S&P 500 · MSCI EM  |  Jan 2004 – Dec 2024  |  1,008 obs."),
    ("Finding",  "KOSPI P/B averages –0.177× below TOPIX (t = –3.23) and –0.601× below MSCI EM (t = –10.30)"),
    ("Approach", "Japan's three governance reforms as policy benchmarks; event-window, panel OLS, synthetic control"),
    ("Caveat",   "N = 4 country clusters; saturated event-window design — descriptive, not causally identified"),
]:
    p = meta.add_paragraph(); p.space_before = Pt(8)
    rl = p.add_run(); rl.text = lbl + ":  "
    rl.font.name = BODY_FONT; rl.font.size = Pt(14); rl.font.bold = True; rl.font.color.rgb = BLACK
    rv = p.add_run(); rv.text = val
    rv.font.name = BODY_FONT; rv.font.size = Pt(14); rv.font.color.rgb = BLACK

br = sl.shapes.add_shape(1, 0, Inches(7.38), W, Inches(0.05))
br.fill.solid(); br.fill.fore_color.rgb = BLACK; br.line.fill.background()


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — The Korea Discount
# ════════════════════════════════════════════════════════════════════════════
sl = slide()
title_bar(sl, "The Korea Discount",
          sub="KOSPI equities have persistently underperformed comparable benchmarks for over two decades")

tf = body(sl, top=Inches(1.35))
para(tf, "Magnitude of the discount — 2004–2024 average P/B ratios", size=16, bold=True)

COLS_L = [Inches(0.55), Inches(3.5), Inches(6.3), Inches(8.7), Inches(10.8)]
COLS_W = [Inches(2.9),  Inches(2.7), Inches(2.35),Inches(2.05),Inches(2.2)]
rows = [
    ("Benchmark", "Mean gap", "t-stat", "95% CI", "KOSPI sub-period"),
    ("vs. TOPIX",   "–0.177×", "–3.23",  "[–0.284×, –0.069×]", ""),
    ("vs. MSCI EM", "–0.601×", "–10.30", "[–0.716×, –0.486×]", ""),
    ("vs. S&P 500", "–1.90×",  "—",      "—",                  ""),
]
sub_notes = ["", "Pre-reform: 1.36×  →  Reform era: 1.02×  →  Post-2023: 0.93×", "", ""]

r_top = Inches(1.95); r_h = Inches(0.47)
for ri, row in enumerate(rows):
    table_row(sl, row, COLS_L, COLS_W, r_top + ri * r_h, r_h,
              is_header=(ri == 0), shade=(ri % 2 == 0))

# sub-period note inline
note_tf = txb(sl, COLS_L[0], r_top + r_h + Pt(4), Inches(12.2), Pt(22))
note_p = note_tf.paragraphs[0]; note_r = note_p.add_run()
note_r.text = "KOSPI trajectory: 1.36× (2004–2013)  →  1.02× (2014–2022)  →  0.93× (2023–2024)  |  TOPIX rebounded to 1.36× post-2023 TSE reform"
note_r.font.name = BODY_FONT; note_r.font.size = Pt(12)
note_r.font.italic = True; note_r.font.color.rgb = MGRAY

hline(sl, Inches(3.5))

para(tf, "\nThree structural channels", size=16, bold=True, space=4)
for ch in [
    "Chaebol opacity & agency costs — control-to-cash-flow ratio ≈ 2:1; tunneling; 60% inheritance tax incentive to suppress share prices",
    "Minority-shareholder recourse deficit — civil-law tradition, weak derivative suit rights, incomplete Stewardship Code",
    "Geopolitical risk premium — six North Korean nuclear tests; priced as a structural level discount, not event spikes",
]:
    bullet(tf, ch, size=14)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — Figure 1: P/B Comparison
# ════════════════════════════════════════════════════════════════════════════
sl = slide()
title_bar(sl, "Figure 1: Index-Level P/B Ratios, 2004–2024",
          sub="Dashed verticals mark Japan's three reform dates; Korea discount persists and widens post-2023")

img(sl, "figure1_pb_comparison.png",
    Inches(0.3), Inches(1.35), Inches(12.7), Inches(5.9))


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — Japan as Policy Benchmark
# ════════════════════════════════════════════════════════════════════════════
sl = slide()
title_bar(sl, "Japan's Governance Reform Program",
          sub="Three staged interventions — used as policy benchmarks, not clean natural experiments")

reforms = [
    ("Feb 2014", "Stewardship Code",
     "Required institutional investors to adopt engagement policies and disclose voting records. "
     "Accelerated keiretsu cross-shareholding unwinding (Miyajima 2023). High compliance under comply-or-explain."),
    ("Jun 2015", "Corporate Governance Code",
     "Required all TSE-listed firms to comply or explain against governance principles. "
     "Mandated ≥ 2 independent directors. By 2017, >90% of JPX-Nikkei 400 complied."),
    ("Mar 2023", "TSE P/B Reform",
     "Firms below 1.0× book value required to disclose capital efficiency plans. "
     "TOPIX P/B rose from ≈1.0× to ≈1.36×. KOSPI remained at ≈0.93× — gap widened. "
     "Korea's Value-Up Program (2024) mirrors this but remains voluntary."),
]

col_l = [Inches(0.55), Inches(1.9), Inches(4.1)]
col_w = [Inches(1.3),  Inches(2.15),Inches(9.0)]
r_top = Inches(1.35); r_h = Inches(1.65)

for ri, (date, name, desc) in enumerate(reforms):
    for ci, (cl, cw, fg) in enumerate(zip(col_l, col_w,
                                          [BLACK, DGRAY, LGRAY if ri % 2 == 0 else WHITE])):
        bg = sl.shapes.add_shape(1, cl, r_top + ri * r_h, cw, r_h)
        bg.fill.solid(); bg.fill.fore_color.rgb = fg; bg.line.fill.background()

    for cl, cw, txt, fc, fsz, fbd in [
        (col_l[0], col_w[0], date, WHITE, 16, True),
        (col_l[1], col_w[1], name, WHITE, 14, True),
        (col_l[2], col_w[2], desc, BLACK, 13, False),
    ]:
        tf2 = txb(sl, cl + Pt(5), r_top + ri * r_h + Pt(7), cw - Pt(10), r_h - Pt(10))
        p2 = tf2.paragraphs[0]; r2 = p2.add_run(); r2.text = txt
        r2.font.name = BODY_FONT; r2.font.size = Pt(fsz)
        r2.font.bold = fbd; r2.font.color.rgb = fc

note = txb(sl, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.55))
note.paragraphs[0].add_run().text = (
    "Korea comparison: adopted Stewardship Code Dec 2016 (weaker enforcement) and Value-Up Program 2024 (voluntary) — "
    "a partial, not strictly never-treated, comparison unit.  "
    "All reform dates locked in config.py before any data transformation.")
note.paragraphs[0].runs[0].font.size = Pt(11)
note.paragraphs[0].runs[0].font.italic = True
note.paragraphs[0].runs[0].font.color.rgb = MGRAY


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — Data & Empirical Strategy
# ════════════════════════════════════════════════════════════════════════════
sl = slide()
title_bar(sl, "Data & Empirical Strategy")

# left col
ltf = txb(sl, Inches(0.55), Inches(1.35), Inches(5.8), Inches(5.85))
para(ltf, "Data", size=17, bold=True)
for d in [
    "Bloomberg PX_TO_BOOK_RATIO: monthly index P/B",
    "KOSPI · TOPIX · S&P 500 · MSCI EM",
    "Jan 2004 – Dec 2024  (1,008 country-months)",
    "P/B preferred to P/E: less cyclically volatile; no sign-change in recessions; directly targeted by TSE reform",
    "GPR index (Caldara–Iacoviello) + GPRNK sub-index for North Korea escalation",
    "No missing data across all four series",
]:
    bullet(ltf, d, size=13)

vl = sl.shapes.add_shape(1, Inches(6.55), Inches(1.35), Pt(1), Inches(5.85))
vl.fill.solid(); vl.fill.fore_color.rgb = LGRAY; vl.line.fill.background()

# right col
rtf = txb(sl, Inches(6.75), Inches(1.35), Inches(6.1), Inches(5.85))
para(rtf, "Three Empirical Designs", size=17, bold=True)

for name, desc in [
    ("Event-window analysis  [primary]",
     "Linear pre-event trend extracted over [–36, –1] months; CAR = cumulative abnormal KOSPI–TOPIX spread "
     "over [–12, +24]. Saturated design → no SEs. Results are descriptive only."),
    ("Panel OLS  [corroborating]",
     "Two-way FE (country + time); reform × Japan interaction terms. "
     "Wild-bootstrap p-values (999 Rademacher draws, N = 4 clusters). "
     "Results: p = 0.625 / 0.375 / 0.500 — no robust causal evidence."),
    ("Synthetic control  [2023 only]",
     "Convex combination of regional EM donor markets on demeaned P/B. "
     "Full weight to MSCI EM Asia (w = 1.0); pre-treatment RMSPE = 0.1451. "
     "Post-reform gap fades rather than widens — no sustained idiosyncratic premium."),
]:
    para(rtf, name, size=14, bold=True, space=8)
    para(rtf, desc, size=13, space=2, indent=1)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — Event-Window Results + Figure 2
# ════════════════════════════════════════════════════════════════════════════
sl = slide()
title_bar(sl, "Results: Event-Window Analysis",
          sub="Cumulative abnormal KOSPI–TOPIX P/B spread; positive = discount narrowed vs. pre-event trend")

img(sl, "figure2_event_study.png",
    Inches(0.3), Inches(1.35), Inches(7.5), Inches(5.85))

rtf = txb(sl, Inches(8.0), Inches(1.35), Inches(5.0), Inches(5.85))
para(rtf, "Three cohort CARs at τ = +24", size=15, bold=True)

for cohort, car, note in [
    ("2014 Stewardship Code",       "CAR ≈ +4.29",
     "Spread rose vs. pre-event trend. Overlaps 2015 window — not independent."),
    ("2015 Corp. Governance Code",  "CAR ≈ +7.83",
     "Largest cumulative movement. Overlaps 2014 window."),
    ("2023 TSE P/B Reform",         "CAR ≈ –6.48",
     "Korea–Japan gap deepened relative to pre-event trend after reform. "
     "Clearest and most policy-relevant cohort."),
]:
    para(rtf, cohort + "  |  " + car, size=14, bold=True, space=10)
    para(rtf, note, size=12, space=1, indent=1, color=MGRAY)

para(rtf,
    "\nCritical caveat: design is saturated; no confidence intervals; "
    "2014 and 2015 cohorts overlap. These are descriptive spread movements, "
    "not causal treatment effects.",
    size=12, italic=True, color=MGRAY, space=10)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — Panel OLS + Synthetic Control
# ════════════════════════════════════════════════════════════════════════════
sl = slide()
title_bar(sl, "Results: Panel OLS & Synthetic Control")

# OLS (left)
ltf = txb(sl, Inches(0.55), Inches(1.35), Inches(6.0), Inches(5.85))
para(ltf, "Panel OLS — two-way FE + reform × Japan", size=16, bold=True)
para(ltf, "Wild-bootstrap p-values (999 Rademacher draws, N = 4 clusters)", size=12,
     italic=True, color=MGRAY, space=2)

ols_rows = [
    ("Interaction term",                  "Coef.",  "p-value"),
    ("Stewardship Code × Japan",          "+0.09",  "0.625"),
    ("Corp. Governance Code × Japan",     "–0.32",  "0.375"),
    ("TSE P/B Reform × Japan",            "–0.23",  "0.500"),
]
cl = [Inches(0.55), Inches(4.0), Inches(5.15)]
cw = [Inches(3.4),  Inches(1.1), Inches(1.35)]
rt = Inches(2.35); rh = Inches(0.48)

for ri, row in enumerate(ols_rows):
    table_row(sl, row, cl, cw, rt + ri * rh, rh, is_header=(ri == 0), shade=(ri % 2 == 0))

para(ltf, "\n2 of 3 point estimates are negative.\nNone reaches p < 0.05 under wild bootstrap.\nN = 4 clusters → severely limited power.\nInterpret as: no robust causal evidence.",
     size=13, space=6)

# divider
vl = sl.shapes.add_shape(1, Inches(6.7), Inches(1.35), Pt(1), Inches(5.85))
vl.fill.solid(); vl.fill.fore_color.rgb = LGRAY; vl.line.fill.background()

# Synthetic control (right)
rtf = txb(sl, Inches(6.9), Inches(1.35), Inches(6.1), Inches(5.85))
para(rtf, "Synthetic Control — 2023 TSE P/B Reform", size=16, bold=True)
para(rtf, "Donor pool: MSCI EM Asia, MSCI HK, MSCI Taiwan, MSCI China, MSCI India, MSCI Indonesia",
     size=12, italic=True, color=MGRAY, space=2)

for lbl, val in [
    ("Donor weight",        "MSCI EM Asia = 1.0  (single-donor solution)"),
    ("Pre-treatment RMSPE", "0.1451  (below 0.15 threshold)"),
    ("Post-reform gap",     "Modestly positive, fading  (avg. change ≈ –0.0028/month)"),
    ("Interpretation",      "Japan's P/B lift broadly matched its synthetic EM Asia peer. "
                            "No sustained idiosyncratic governance premium after 2023."),
]:
    p = rtf.add_paragraph(); p.space_before = Pt(9)
    r1 = p.add_run(); r1.text = lbl + ":  "
    r1.font.name = BODY_FONT; r1.font.size = Pt(14); r1.font.bold = True; r1.font.color.rgb = BLACK
    r2 = p.add_run(); r2.text = val
    r2.font.name = BODY_FONT; r2.font.size = Pt(14); r2.font.color.rgb = BLACK

para(rtf,
    "\nCaveat: TOPIX absolute P/B did rise, but may reflect yen depreciation and "
    "regional EM momentum rather than a purely idiosyncratic governance premium.",
    size=12, italic=True, color=MGRAY, space=8)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — Figures: Geo Risk + Synth Gap
# ════════════════════════════════════════════════════════════════════════════
sl = slide()
title_bar(sl, "Preliminary Results: Additional Figures",
          sub="Left: GPR escalation insignificant (p = 0.40) — risk priced as level discount   |   Right: Synthetic control gap fades post-2023")

img(sl, "figure3_geo_risk.png",
    Inches(0.3), Inches(1.35), Inches(6.3), Inches(5.85))

img(sl, "figure_synth_gap.png",
    Inches(6.9), Inches(1.35), Inches(6.3), Inches(5.85))


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — Counterfactual Projection
# ════════════════════════════════════════════════════════════════════════════
sl = slide()
title_bar(sl, "Illustrative Counterfactual Projection",
          sub="KOSPI P/B under a Korea governance reform scenario — stress test, not a forecast")

img(sl, "figure4_counterfactual_projection.png",
    Inches(0.3), Inches(1.35), Inches(7.5), Inches(5.85))

rtf = txb(sl, Inches(8.0), Inches(1.35), Inches(5.0), Inches(5.85))
para(rtf, "Construction", size=16, bold=True)
for d in [
    "Base: KOSPI P/B ≈ 0.93× (Dec 2024)",
    "Applied rate: avg. monthly change in Japan's synthetic-control gap over first 18 post-reform months (–0.0028/month)",
    "Uncertainty band: ± RMSPE = ± 0.1451",
]:
    bullet(rtf, d, size=13)

para(rtf, "\nResult", size=16, bold=True, space=10)
para(rtf,
    "Because the Japan post-2023 relative gap change is negative, the mechanical projection "
    "does not close the Korea Discount. This is a stress test, not a prediction of policy success.",
    size=13, space=2)

para(rtf, "\nPolicy levers (if reforms are implemented)", size=15, bold=True, space=10)
for rec in [
    "FSC mandatory P/B disclosure (comply-or-explain, not voluntary)",
    "KRX ≥ 50% independent board composition for large-cap issuers",
    "Korean Stewardship Code: mandatory voting records + NPS engagement division",
]:
    bullet(rtf, rec, size=13)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — Conclusion
# ════════════════════════════════════════════════════════════════════════════
sl = slide()
title_bar(sl, "Conclusion")

tf = body(sl, top=Inches(1.35))

for heading, bullets in [
    ("What we document",
     [
         "KOSPI P/B averages –0.177× below TOPIX (t = –3.23) and –0.601× below MSCI EM (t = –10.30)",
         "Discount widened post-2023: KOSPI ≈ 0.93× vs. TOPIX ≈ 1.36×",
         "Three channels: chaebol opacity · minority recourse deficit · geopolitical risk premium",
     ]),
    ("Empirical findings",
     [
         "Event-window CARs: large spread movements around reform dates, but no valid inference (saturated design, overlapping cohorts)",
         "Panel OLS: p-values 0.625 / 0.375 / 0.500 — no robust causal evidence of a Japan governance re-rating",
         "Synthetic control: full weight to MSCI EM Asia (RMSPE 0.1451); post-reform gap fades, not widens",
         "Geopolitical risk coefficient: –0.02 (p = 0.40) — priced as a structural level discount",
     ]),
    ("Bottom line",
     [
         "The Korea Discount is persistent and statistically significant",
         "Japan's experience provides useful policy benchmarks but does not identify the valuation gain "
           "Korea would realize from adopting Japan-style reforms",
         "Future work: firm-level chaebol governance panel; Korea Value-Up Program event study (2024–2026+)",
     ]),
]:
    para(tf, heading, size=16, bold=True, space=8)
    for b in bullets:
        bullet(tf, b, size=13)


# ════════════════════════════════════════════════════════════════════════════
#  Save
# ════════════════════════════════════════════════════════════════════════════
out = "/Users/dandan/Desktop/Projects/kor-discount/output/korea_discount_slides_short.pptx"
prs.save(out)
print(f"Saved → {out}")
