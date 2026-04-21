"""
Generate a structured PowerPoint presentation for the Korea Discount paper.
Structure: Introduction → Literature Review → Methodology → Results → Discussion
Content drawn directly from main_v2.tex.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── colour palette ──────────────────────────────────────────────────────────
BLACK  = RGBColor(0x00, 0x00, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xD9, 0xD9, 0xD9)
MGRAY  = RGBColor(0x59, 0x59, 0x59)
DGRAY  = RGBColor(0x1A, 0x1A, 0x1A)

TITLE_FONT = "Calibri"
BODY_FONT  = "Calibri"

# ── slide dimensions (widescreen 16:9) ──────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

FIGURES = "/Users/dandan/Desktop/Projects/kor-discount/output/figures/png"


# ═══════════════════════════════════════════════════════════════════════════
#  Helpers
# ═══════════════════════════════════════════════════════════════════════════

def add_slide():
    return prs.slides.add_slide(BLANK)


def txb(slide, left, top, width, height):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    return tf


def title_bar(slide, text, subtitle=None):
    """Thin black top rule, title text, optional subtitle, bottom rule."""
    rule = slide.shapes.add_shape(1, 0, 0, W, Inches(0.07))
    rule.fill.solid(); rule.fill.fore_color.rgb = BLACK; rule.line.fill.background()

    tf = txb(slide, Inches(0.55), Inches(0.15), Inches(12.2), Inches(0.75))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.name = TITLE_FONT; r.font.bold = True
    r.font.size = Pt(28); r.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        tf2 = txb(slide, Inches(0.55), Inches(0.85), Inches(12.2), Inches(0.4))
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.name = BODY_FONT; r2.font.size = Pt(14)
        r2.font.color.rgb = MGRAY; p2.alignment = PP_ALIGN.LEFT

    br = slide.shapes.add_shape(1, 0, Inches(7.38), W, Inches(0.05))
    br.fill.solid(); br.fill.fore_color.rgb = BLACK; br.line.fill.background()


def body_tf(slide, top=Inches(1.35), left=Inches(0.55),
            width=Inches(12.2), height=Inches(5.85)):
    return txb(slide, left, top, width, height)


def add_para(tf, text, size=17, bold=False, indent=0, color=BLACK,
             space_before=0, italic=False, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.space_before = Pt(space_before)
    if indent:
        p.level = indent
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = BODY_FONT
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def section_divider(label, sub=None):
    """Full-bleed black slide with white section label."""
    sl = add_slide()
    bg = sl.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BLACK; bg.line.fill.background()

    top = Inches(2.9) if not sub else Inches(2.6)
    tf = txb(sl, Inches(1.5), top, Inches(10), Inches(1.2))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = label
    r.font.name = TITLE_FONT; r.font.size = Pt(40)
    r.font.bold = True; r.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if sub:
        tf2 = txb(sl, Inches(1.5), Inches(3.85), Inches(10), Inches(0.8))
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run(); r2.text = sub
        r2.font.name = BODY_FONT; r2.font.size = Pt(18)
        r2.font.bold = False; r2.font.color.rgb = LGRAY
        p2.alignment = PP_ALIGN.CENTER

    return sl


def hline(slide, top, color=LGRAY, left=Inches(0.55), width=Inches(12.2)):
    ln = slide.shapes.add_shape(1, left, top, width, Pt(1))
    ln.fill.solid(); ln.fill.fore_color.rgb = color; ln.line.fill.background()


def add_figure(slide, img_path, left, top, width, height):
    slide.shapes.add_picture(img_path, left, top, width, height)


def two_col_split(slide, left_title, left_items, right_title, right_items,
                  top=Inches(1.35), col_w=Inches(5.9), gap=Inches(0.5),
                  size=14, header_size=17):
    """Helper for two-column text layout."""
    left_l  = Inches(0.55)
    right_l = Inches(0.55) + col_w + gap

    ltf = txb(slide, left_l, top, col_w, Inches(5.85))
    add_para(ltf, left_title, size=header_size, bold=True)
    for item in left_items:
        add_para(ltf, item, size=size, space_before=3)

    vl = slide.shapes.add_shape(1, right_l - gap / 2, top, Pt(1), Inches(5.85))
    vl.fill.solid(); vl.fill.fore_color.rgb = LGRAY; vl.line.fill.background()

    rtf = txb(slide, right_l, top, col_w, Inches(5.85))
    add_para(rtf, right_title, size=header_size, bold=True)
    for item in right_items:
        add_para(rtf, item, size=size, space_before=3)

    return ltf, rtf


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 1 – Title
# ═══════════════════════════════════════════════════════════════════════════
sl = add_slide()

bar = sl.shapes.add_shape(1, 0, 0, W, Inches(3.1))
bar.fill.solid(); bar.fill.fore_color.rgb = BLACK; bar.line.fill.background()

tf = txb(sl, Inches(0.7), Inches(0.4), Inches(11.9), Inches(2.4))
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Corporate Governance Reform and the Korea Discount"
r.font.name = TITLE_FONT; r.font.size = Pt(34); r.font.bold = True
r.font.color.rgb = WHITE; p.alignment = PP_ALIGN.LEFT

p2 = tf.add_paragraph(); p2.space_before = Pt(6)
r2 = p2.add_run()
r2.text = "Lessons from Japan's Reform Experience"
r2.font.name = TITLE_FONT; r2.font.size = Pt(21)
r2.font.bold = False; r2.font.color.rgb = LGRAY
p2.alignment = PP_ALIGN.LEFT

meta = txb(sl, Inches(0.7), Inches(3.25), Inches(11.9), Inches(3.8))
items = [
    ("Setting",
     "Monthly index-level P/B ratios for KOSPI, TOPIX, S&P 500, MSCI EM  ·  2004–2024  ·  1,008 country-months"),
    ("Core finding",
     "KOSPI P/B averages –0.177× below TOPIX (t = –3.23) and –0.601× below MSCI EM (t = –10.30): "
     "a persistent, statistically significant Korea Discount"),
    ("Approach",
     "Japan's three staged governance reforms used as policy benchmarks; "
     "descriptive event-window analysis, panel OLS (wild-bootstrap), and synthetic control"),
    ("Caveat",
     "N = 4 country clusters; saturated event-window design; evidence is descriptive, not causal"),
]
for lbl, val in items:
    p = meta.add_paragraph(); p.space_before = Pt(7)
    rl = p.add_run(); rl.text = lbl + ":  "
    rl.font.name = BODY_FONT; rl.font.size = Pt(14); rl.font.bold = True
    rl.font.color.rgb = BLACK
    rv = p.add_run(); rv.text = val
    rv.font.name = BODY_FONT; rv.font.size = Pt(14); rv.font.bold = False
    rv.font.color.rgb = BLACK

br = sl.shapes.add_shape(1, 0, Inches(7.38), W, Inches(0.05))
br.fill.solid(); br.fill.fore_color.rgb = BLACK; br.line.fill.background()


# ═══════════════════════════════════════════════════════════════════════════
#  SECTION: INTRODUCTION
# ═══════════════════════════════════════════════════════════════════════════
section_divider("Introduction", "The Korea Discount Puzzle")


# ─── Slide: The Korea Discount ───────────────────────────────────────────
sl = add_slide()
title_bar(sl, "The Korea Discount",
          subtitle="KOSPI equities trade at a persistent structural discount to developed-market and EM peers")

tf = body_tf(sl, top=Inches(1.35))

add_para(tf, "The puzzle", size=18, bold=True)
add_para(tf,
    "Despite Korea's G20 membership, high-income fundamentals, and a well-developed export sector, "
    "KOSPI index-level price-to-book ratios have persistently trailed comparable benchmarks for over two decades.",
    size=15, space_before=3)

hline(sl, Inches(2.35))

add_para(tf, "Magnitude of the discount  (2004–2024 average P/B ratios)", size=18, bold=True, space_before=8)

rows = [
    ("Benchmark",  "KOSPI mean P/B", "Benchmark mean P/B", "Mean gap",  "t-statistic",   "95% CI"),
    ("TOPIX",      "1.18×",          "1.35×",               "–0.177×",  "–3.23",          "[–0.284×, –0.069×]"),
    ("MSCI EM",    "1.18×",          "1.78×",               "–0.601×",  "–10.30",         "[–0.716×, –0.486×]"),
    ("S&P 500",    "1.18×",          "3.08×",               "–1.90×",   "—",              "—"),
]

col_w = [Inches(2.0), Inches(1.7), Inches(2.0), Inches(1.5), Inches(1.6), Inches(2.5)]
col_l = [Inches(0.55), Inches(2.6), Inches(4.35), Inches(6.4), Inches(7.95), Inches(9.6)]
r_top = Inches(3.0)
r_h   = Inches(0.43)

for ri, row in enumerate(rows):
    is_hdr = (ri == 0)
    for ci, (cell, cw, cl) in enumerate(zip(row, col_w, col_l)):
        bg = sl.shapes.add_shape(1, cl, r_top + ri * r_h, cw, r_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BLACK if is_hdr else (LGRAY if ri % 2 == 0 else WHITE)
        bg.line.fill.background()
        ctf = txb(sl, cl + Pt(4), r_top + ri * r_h + Pt(4), cw - Pt(8), r_h - Pt(5))
        cp = ctf.paragraphs[0]
        cr = cp.add_run(); cr.text = cell
        cr.font.name = BODY_FONT; cr.font.size = Pt(13)
        cr.font.bold = is_hdr
        cr.font.color.rgb = WHITE if is_hdr else BLACK
        cp.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT

add_para(tf,
    "\nSub-period: KOSPI P/B fell from 1.36× (2004–2013) → 1.02× (2014–2022) → 0.93× (2023–2024), "
    "while TOPIX rebounded to 1.36× after the 2023 TSE P/B reform.",
    size=13, italic=True, color=MGRAY, space_before=4)


# ─── Slide: Research Contributions ──────────────────────────────────────
sl = add_slide()
title_bar(sl, "Research Contributions",
          subtitle="Three contributions to the Korea Discount literature")

tf = body_tf(sl, top=Inches(1.35))

contribs = [
    ("Contribution 1:  Descriptive event-window analysis",
     "First panel study using all three Japan governance reform dates jointly as descriptive event windows "
     "for the KOSPI–TOPIX valuation spread. Documents timing and magnitude of spread movements. "
     "Explicitly notes that the saturated design and overlapping 2014–2015 windows do not deliver "
     "valid conventional inference."),
    ("Contribution 2:  Panel OLS with honest inference",
     "Two-way fixed-effects panel estimates with small-cluster wild-bootstrap p-values (N = 4 countries, "
     "999 Rademacher draws). Shows that reform-by-Japan interaction terms are too noisy — "
     "p-values of 0.625, 0.375, and 0.500 — to sustain strong directional claims."),
    ("Contribution 3:  Synthetic control diagnostic",
     "Applies Abadie (2010) synthetic control to the 2023 TSE P/B reform. Transparently reports the "
     "single-donor weight structure (MSCI EM Asia, weight = 1.0) and pre-treatment RMSPE = 0.1451. "
     "Limits interpretation accordingly: the post-reform gap does not isolate a purely idiosyncratic "
     "Japan governance premium after accounting for regional equity momentum."),
]

for title, body in contribs:
    add_para(tf, title, size=16, bold=True, space_before=8)
    add_para(tf, body, size=14, space_before=2, indent=1)

add_para(tf, "\nRoad map:  Lit Review → Methodology → Results → Discussion",
         size=13, italic=True, color=MGRAY, space_before=8)


# ═══════════════════════════════════════════════════════════════════════════
#  SECTION: LITERATURE REVIEW
# ═══════════════════════════════════════════════════════════════════════════
section_divider("Literature Review", "Chaebol Structure · Japan Reforms · Governance Theory · Geopolitical Risk")


# ─── Slide: Korea Discount Literature & Chaebol Structure ───────────────
sl = add_slide()
title_bar(sl, "The Korea Discount: Prior Evidence",
          subtitle="Chaebol ownership architecture and principal–principal agency costs")

tf = body_tf(sl, top=Inches(1.35))

add_para(tf, "Chaebol structure (Claessens et al. 2000; Baek et al. 2004)", size=17, bold=True)
items = [
    "Family-controlled conglomerates dominate KOSPI; control-to-cash-flow ratio ≈ 2:1 (highest in East Asia)",
    "Principal–principal agency problem: controlling family vs. minority shareholders (distinct from Jensen-Meckling)",
    "Tunneling: intra-group related-party transactions that extract private benefits at minority expense",
    "Families incentivized to keep share prices depressed (Korea inheritance tax up to 60%, applied to market value)",
    "\"Double-counting\" in index-level valuation inflates book denominator → mechanically lower P/B",
]
for item in items:
    add_para(tf, "•  " + item, size=14, space_before=2, indent=1)

hline(sl, Inches(3.65))

add_para(tf, "Key empirical findings", size=17, bold=True, space_before=8)
findings = [
    "Baek et al. (2004): Firms with larger control–cash-flow wedge saw larger equity declines in 1997 crisis",
    "Black et al. (2006): Higher KCGI governance index → significantly higher Tobin's q (panel OLS, firm-level)",
    "KCMI (2023): Four channels — governance structure, low shareholder returns, market accessibility, inheritance tax",
    "This paper: operates at country-level index — fills firm-level literature's gap on macro reform shocks, "
      "at the cost of N = 1 treated country",
]
for item in findings:
    add_para(tf, "•  " + item, size=14, space_before=2, indent=1)


# ─── Slide: Japan's Three Governance Reforms ────────────────────────────
sl = add_slide()
title_bar(sl, "Japan's Governance Reform Program",
          subtitle="Three staged interventions — policy benchmarks, not clean natural experiments")

reforms = [
    ("Feb 2014",  "Japan Stewardship Code",
     "FSA required institutional investors to adopt engagement policies and disclose voting records. "
     "Comply-or-explain format. Miyajima (2023): accelerated unwinding of keiretsu cross-shareholdings; "
     "firms with larger cross-holding reductions showed commensurately larger Tobin's q improvements."),
    ("Jun 2015",  "Corporate Governance Code",
     "TSE + FSA required all listed firms to comply or explain against governance principles. "
     "Mandated ≥ 2 independent outside directors. By 2017, >90% of JPX-Nikkei 400 had ≥ 2 independent directors. "
     "Eberhart (2012): board independence → operating performance and equity valuation improvements."),
    ("Mar 2023",  "TSE P/B Reform",
     "Instructed companies trading below 1.0× book value to disclose capital efficiency plans. "
     "TOPIX P/B rose from ≈ 1.0× (early 2023) to ≈ 1.36× (2023–2024). "
     "KOSPI remained at ≈ 0.93× → Korea–Japan P/B gap widened. "
     "Korea's Value-Up Program (2024) is modelled explicitly on this reform but remains voluntary."),
]

col_w = [Inches(1.5), Inches(2.5), Inches(8.3)]
col_l = [Inches(0.55), Inches(2.1), Inches(4.65)]
r_top = Inches(1.35)
r_h   = Inches(1.65)

for ri, (date, name, desc) in enumerate(reforms):
    # date cell (black)
    bg0 = sl.shapes.add_shape(1, col_l[0], r_top + ri * r_h, col_w[0], r_h)
    bg0.fill.solid(); bg0.fill.fore_color.rgb = BLACK; bg0.line.fill.background()
    dtf = txb(sl, col_l[0] + Pt(4), r_top + ri * r_h + Pt(8), col_w[0] - Pt(8), r_h - Pt(12))
    dp = dtf.paragraphs[0]
    dr = dp.add_run(); dr.text = date
    dr.font.name = BODY_FONT; dr.font.size = Pt(16); dr.font.bold = True
    dr.font.color.rgb = WHITE; dp.alignment = PP_ALIGN.CENTER

    # name cell (dark gray)
    bg1 = sl.shapes.add_shape(1, col_l[1], r_top + ri * r_h, col_w[1], r_h)
    bg1.fill.solid(); bg1.fill.fore_color.rgb = DGRAY; bg1.line.fill.background()
    ntf = txb(sl, col_l[1] + Pt(6), r_top + ri * r_h + Pt(8), col_w[1] - Pt(12), r_h - Pt(12))
    np_ = ntf.paragraphs[0]
    nr = np_.add_run(); nr.text = name
    nr.font.name = BODY_FONT; nr.font.size = Pt(14); nr.font.bold = True
    nr.font.color.rgb = WHITE

    # desc cell (alternating)
    bg2 = sl.shapes.add_shape(1, col_l[2], r_top + ri * r_h, col_w[2], r_h)
    bg2.fill.solid(); bg2.fill.fore_color.rgb = LGRAY if ri % 2 == 0 else WHITE
    bg2.line.fill.background()
    dtf2 = txb(sl, col_l[2] + Pt(6), r_top + ri * r_h + Pt(6), col_w[2] - Pt(12), r_h - Pt(8))
    dp2 = dtf2.paragraphs[0]
    dr2 = dp2.add_run(); dr2.text = desc
    dr2.font.name = BODY_FONT; dr2.font.size = Pt(12)
    dr2.font.color.rgb = BLACK

note = txb(sl, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.55))
np2 = note.paragraphs[0]
nr2 = np2.add_run()
nr2.text = ("All three reform dates locked in config.py before any data transformation — "
            "strict look-ahead bias firewall. Korea is a partial comparison case: "
            "adopted its own Stewardship Code (Dec 2016) and Value-Up Program (2024).")
nr2.font.name = BODY_FONT; nr2.font.size = Pt(12); nr2.font.italic = True
nr2.font.color.rgb = MGRAY


# ─── Slide: Governance Theory & Geopolitical Risk ───────────────────────
sl = add_slide()
title_bar(sl, "Theoretical Grounding",
          subtitle="Law-and-finance framework · Geopolitical risk channel")

ltf = txb(sl, Inches(0.55), Inches(1.35), Inches(5.85), Inches(5.85))
add_para(ltf, "Governance–Valuation Theory", size=17, bold=True)
theory = [
    "Shleifer & Vishny (1997): ownership concentration resolves manager–shareholder conflict "
    "but intensifies controlling-shareholder vs. minority conflict — exactly the chaebol problem",
    "La Porta et al. (1998, 2002): stronger statutory investor protections → larger, more liquid "
    "capital markets; minority protection causally associated with higher dividends and Tobin's q",
    "Korea's civil-law tradition + chaebol structure places it toward the weak-protection end "
    "of the international distribution",
    "Japan's reform experience demonstrates that de facto governance can be improved through "
    "targeted regulation even within a civil-law system",
    "Gompers et al. (2003): firms with stronger shareholder rights earn significantly higher "
    "abnormal returns (US evidence; corroborating)",
]
for item in theory:
    add_para(ltf, "•  " + item, size=13, space_before=4, indent=1)

vl = sl.shapes.add_shape(1, Inches(6.6), Inches(1.35), Pt(1), Inches(5.85))
vl.fill.solid(); vl.fill.fore_color.rgb = LGRAY; vl.line.fill.background()

rtf = txb(sl, Inches(6.8), Inches(1.35), Inches(6.0), Inches(5.85))
add_para(rtf, "Geopolitical Risk Channel", size=17, bold=True)
geo = [
    "Six North Korean nuclear tests over sample period: Oct 2006, May 2009, Feb 2013, Jan 2016, Sep 2016, Sep 2017",
    "2017 test (ICBM range) coincided with largest single-period widening of Korea–TOPIX P/B spread",
    "Caldara–Iacoviello GPR index and Korea-specific GPRNK sub-index (IMF 2021)",
    "IMF (2021): North Korean provocations primarily affect systematic risk / market beta — "
    "a level effect on required return — not cumulative permanent discount widening",
    "Implication: geopolitical premium is embedded in the structural discount level, "
    "not generating time-series spikes around each escalation event",
    "Cannot be resolved by domestic governance reform alone",
]
for item in geo:
    add_para(rtf, "•  " + item, size=13, space_before=4, indent=1)


# ═══════════════════════════════════════════════════════════════════════════
#  SECTION: METHODOLOGY
# ═══════════════════════════════════════════════════════════════════════════
section_divider("Methodology", "Data · Causal Channels · Empirical Strategy")


# ─── Slide: Data Sources ────────────────────────────────────────────────
sl = add_slide()
title_bar(sl, "Data Sources & Variable Construction",
          subtitle="Monthly index-level P/B ratios from Bloomberg, 2004–2024")

tf = body_tf(sl, top=Inches(1.35))

add_para(tf, "Primary panel", size=17, bold=True)
data_items = [
    "Bloomberg Terminal  PX_TO_BOOK_RATIO: aggregate market cap ÷ aggregate book value for all index constituents",
    "Four indices: KOSPI Composite (KOSPI Index)  ·  TOPIX (TPX Index)  ·  S&P 500 (SPX Index)  ·  MSCI EM (MXEF Index)",
    "Primary sample: Jan 2004 – Dec 2024  →  252 months per series  →  1,008 country-month observations",
    "Processed panel stored as data/processed/panel.parquet with schema (date, country, pb, pe, fx_rate)",
    "All Bloomberg raw CSV exports version-controlled in data/raw/ with MANIFEST.md (vintage: 2026-04-20)",
    "Geopolitical Risk: Caldara–Iacoviello GPR index + GPRNK sub-index  ·  monthly, no gaps",
]
for item in data_items:
    add_para(tf, "•  " + item, size=14, space_before=3, indent=1)

hline(sl, Inches(4.05))

add_para(tf, "Why P/B over P/E?", size=17, bold=True, space_before=8)
pb_items = [
    "(i)  Book value is considerably less cyclically volatile → more stable measure of structural discount",
    "(ii)  Book value does not go negative for index aggregates (avoids sign-change problem in recessions)",
    "(iii) The TSE P/B reform and Korea's Value-Up Program are both explicitly framed in P/B terms",
]
for item in pb_items:
    add_para(tf, item, size=14, space_before=3, indent=1)

add_para(tf,
    "\nSurvivorship bias: index-level P/B computed over all constituents as of each month-end, "
    "including firms that subsequently exit. Korea Discount persisted despite constituent turnover.",
    size=13, italic=True, color=MGRAY, space_before=4)


# ─── Slide: Three Structural Channels ───────────────────────────────────
sl = add_slide()
title_bar(sl, "Three Structural Channels",
          subtitle="Compounding mechanisms that explain the Korea Discount")

tf = body_tf(sl, top=Inches(1.35))

channels = [
    ("1.  Chaebol Opacity and Agency Costs",
     [
         "Control-to-cash-flow ratio ≈ 2:1 → intragroup tunneling; group accounting obscures entity economics",
         "Controlling-family incentives to suppress equity prices during generational wealth transfers (60% inheritance tax)",
         "Double-counting in KOSPI book value inflates denominator → mechanically lower P/B",
     ]),
    ("2.  Minority Shareholder Recourse Deficit",
     [
         "High ownership thresholds for derivative suits; lax board independence pre-2020 Commercial Act reforms",
         "Related-party transaction disclosure thresholds permitted large intragroup transactions without shareholder approval",
         "Korea Stewardship Code (Dec 2016) widely assessed as weaker than Japan's 2014 equivalent",
         "Commercial Act fiduciary duty amendments for minority shareholders proposed only 2025",
     ]),
    ("3.  Geopolitical Risk Premium",
     [
         "Tail risk of armed conflict warrants a discount to intrinsic value not shared by TOPIX, S&P 500, or MSCI EM",
         "Priced continuously as a level discount rather than reacting discretely to individual escalation events",
         "Cannot be eliminated through domestic governance reform alone",
     ]),
]

for ch_title, bullets in channels:
    add_para(tf, ch_title, size=16, bold=True, space_before=8)
    for b in bullets:
        add_para(tf, "•  " + b, size=13, space_before=2, indent=1)

add_para(tf,
    "\nThe interaction of these three channels creates a cumulative discount exceeding "
    "what any individual channel would predict in isolation.",
    size=13, italic=True, color=MGRAY, space_before=6)


# ─── Slide: Empirical Strategy ──────────────────────────────────────────
sl = add_slide()
title_bar(sl, "Empirical Strategy",
          subtitle="Three designs with different evidentiary weight — all reform dates locked before data transformation")

tf = body_tf(sl, top=Inches(1.35))

designs = [
    ("Descriptive Event-Window Analysis  [primary description]",
     "For each reform date k ∈ {2014, 2015, 2023}, construct monthly spread S_t = PB_KOSPI – PB_TOPIX. "
     "Estimate linear pre-event trend in S_t over [–36, –1] months; define abnormal spread movement as "
     "observed spread minus fitted pre-event trend. Report cumulative abnormal spread (CAR) over [–12, +24]. "
     "Design is saturated at cohort-by-relative-month level → conventional SEs undefined → results are descriptive only. "
     "2014 and 2015 windows overlap — not independent reform experiments."),
    ("Panel OLS  [corroborating]",
     "Two-way FE (country + time) specification: Y_ct = α_c + α_t + Σ_k β_k · D_kt · 1[c = Japan] + ε_ct. "
     "Reform-by-Japan interaction terms β_k estimated via linearmodels.PanelOLS in Python (within-demeaning). "
     "Inference: wild-bootstrap with 999 Rademacher draws, clustered by country (Cameron et al. 2008). "
     "N = 4 country clusters → p-values are conservative but size-correct. "
     "Results: p-values of 0.625, 0.375, 0.500 — not robust evidence of positive re-rating."),
    ("Synthetic Control  [2023 only]",
     "Abadie (2010) convex combination of donor markets: MSCI EM, MSCI EM Asia, MSCI HK, MSCI Taiwan, "
     "MSCI China, MSCI India, MSCI Indonesia. Optimization on demeaned P/B resolves convex-hull boundary. "
     "KOSPI excluded from donor pool. Full weight assigned to MSCI EM Asia (weight = 1.0); "
     "pre-treatment RMSPE = 0.1451. Post-reform gap should not be read as a stand-alone causal estimate."),
]

for title, body in designs:
    add_para(tf, title, size=15, bold=True, space_before=10)
    add_para(tf, body, size=13, space_before=2, indent=1)


# ═══════════════════════════════════════════════════════════════════════════
#  SECTION: RESULTS
# ═══════════════════════════════════════════════════════════════════════════
section_divider("Results", "Event-Windows · Panel OLS · Geopolitical Risk · Synthetic Control")


# ─── Slide: Figure 1 – P/B Comparison ──────────────────────────────────
sl = add_slide()
title_bar(sl, "Figure 1: Index-Level P/B Ratios, 2004–2024",
          subtitle="KOSPI trades at a persistent discount; dashed lines mark Japan's three reform dates")

add_figure(sl,
           f"{FIGURES}/figure1_pb_comparison.png",
           Inches(0.55), Inches(1.35),
           Inches(8.5), Inches(5.85))

note_tf = txb(sl, Inches(9.2), Inches(1.35), Inches(3.9), Inches(5.85))
add_para(note_tf, "Key observations", size=15, bold=True)
obs = [
    "KOSPI P/B consistently below TOPIX, S&P 500, and MSCI EM throughout 2004–2024",
    "Largest single-period widening visible around 2017 North Korean ICBM tests",
    "Post-2023 TSE P/B reform: TOPIX rebounds to 1.36×; KOSPI stays near 0.93×",
    "Korea–Japan gap widens at exactly the point Japan's governance reform was most aggressive",
    "Global financial crisis (2008–09) compresses all four series simultaneously — cross-market consistency check",
]
for o in obs:
    add_para(note_tf, "•  " + o, size=13, space_before=5, indent=1)

add_para(note_tf,
    "\nMean gap vs. TOPIX: –0.177× (t = –3.23)\n"
    "Mean gap vs. MSCI EM: –0.601× (t = –10.30)",
    size=14, bold=True, space_before=10)


# ─── Slide: Event-Window Analysis (text summary) ────────────────────────
sl = add_slide()
title_bar(sl, "Descriptive Event-Window Analysis",
          subtitle="Cumulative abnormal movements in the KOSPI–TOPIX P/B spread around Japan reform dates")

tf = body_tf(sl, top=Inches(1.35))

add_para(tf,
    "Sign convention: positive CAR = KOSPI–TOPIX spread rose relative to pre-event trend "
    "(Korea discount narrowed or reversed); negative CAR = Korea discount deepened.",
    size=13, italic=True, color=MGRAY)

cohorts = [
    ("2014 Stewardship Code",
     "CAR by τ = +24:  ≈ +4.29",
     "Negative pre-event CAR reverses after reform date. Under this sign convention, the spread "
     "rose vs. its pre-event trend — consistent with some narrowing of the Korea discount relative "
     "to its trajectory. However this window overlaps with the 2015 reform window and is not independent."),
    ("2015 Corporate Governance Code",
     "CAR by τ = +24:  ≈ +7.83",
     "Largest cumulative movement of the three cohorts. Pre-event CARs relatively flat. "
     "Heavily overlaps with 2014 window — the two events cannot be treated as independent experiments "
     "given they applied to the same country within 16 months. Standard errors are undefined."),
    ("2023 TSE P/B Reform",
     "CAR by τ = +24:  ≈ –6.48",
     "CARs become persistently negative after March 2023. Under the sign convention, this means "
     "the Korea–Japan spread worsened relative to its pre-event trend: TOPIX re-rated sharply higher "
     "while KOSPI did not. This is the clearest and most recent policy-relevant cohort."),
]

for ch, car, interp in cohorts:
    bg = sl.shapes.add_shape(1, Inches(0.55), Inches(1.7 + cohorts.index((ch, car, interp)) * 1.65),
                              Inches(12.2), Inches(0.32))
    bg.fill.solid(); bg.fill.fore_color.rgb = BLACK; bg.line.fill.background()
    ctf = txb(sl, Inches(0.65), Inches(1.7 + cohorts.index((ch, car, interp)) * 1.65) + Pt(3),
              Inches(12.0), Inches(0.28))
    cp = ctf.paragraphs[0]; cr = cp.add_run(); cr.text = ch + "   " + car
    cr.font.name = BODY_FONT; cr.font.size = Pt(14); cr.font.bold = True; cr.font.color.rgb = WHITE

    btf = txb(sl, Inches(0.75), Inches(1.7 + cohorts.index((ch, car, interp)) * 1.65) + Inches(0.35),
              Inches(12.0), Inches(1.25))
    bp = btf.paragraphs[0]; br2 = bp.add_run(); br2.text = interp
    br2.font.name = BODY_FONT; br2.font.size = Pt(13); br2.font.color.rgb = BLACK

add_para(tf,
    "\nImportant: these are descriptive abnormal spread movements. "
    "Not presented with confidence intervals. Not a causal identification.",
    size=13, italic=True, color=MGRAY, space_before=4)


# ─── Slide: Figure 2 – Event Study ──────────────────────────────────────
sl = add_slide()
title_bar(sl, "Figure 2: Event-Window CARs Around Japan Reform Dates",
          subtitle="Three panels showing cumulative abnormal KOSPI–TOPIX P/B spread; descriptive only, no confidence intervals")

add_figure(sl,
           f"{FIGURES}/figure2_event_study.png",
           Inches(0.3), Inches(1.35),
           Inches(13.0), Inches(5.9))


# ─── Slide: Panel OLS Results ───────────────────────────────────────────
sl = add_slide()
title_bar(sl, "Panel OLS Results",
          subtitle="Two-way FE (country + time); wild-bootstrap p-values in brackets; N = 4 country clusters")

add_para(body_tf(sl, top=Inches(1.35), height=Inches(0.4)),
    "Specification: Y_ct = α_c + α_t + Σ_k β_k · D_kt · 1[c = Japan] + ε_ct",
    size=14, italic=True, color=MGRAY)

rows = [
    ("Reform Interaction",                   "Coefficient",  "Wild-bootstrap p-value"),
    ("Stewardship Code × Japan",             "+0.09",        "0.625"),
    ("Corporate Governance Code × Japan",    "–0.32",        "0.375"),
    ("TSE P/B Reform × Japan",               "–0.23",        "0.500"),
]

col_w = [Inches(6.5), Inches(2.4), Inches(3.0)]
col_l = [Inches(0.55), Inches(7.1), Inches(9.6)]
r_top = Inches(2.1)
r_h   = Inches(0.55)

for ri, row in enumerate(rows):
    is_hdr = (ri == 0)
    for ci, (cell, cw, cl) in enumerate(zip(row, col_w, col_l)):
        bg = sl.shapes.add_shape(1, cl, r_top + ri * r_h, cw, r_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BLACK if is_hdr else (LGRAY if ri % 2 == 0 else WHITE)
        bg.line.fill.background()
        ctf = txb(sl, cl + Pt(6), r_top + ri * r_h + Pt(5), cw - Pt(12), r_h - Pt(6))
        cp = ctf.paragraphs[0]; cr = cp.add_run(); cr.text = cell
        cr.font.name = BODY_FONT; cr.font.size = Pt(15)
        cr.font.bold = is_hdr; cr.font.color.rgb = WHITE if is_hdr else BLACK
        cp.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT

tf2 = body_tf(sl, top=Inches(4.45), height=Inches(2.7))
takeaways = [
    "Two of three point estimates are negative — directionally inconsistent with the simple 'Japan re-rated up' narrative",
    "None of the three reform-by-Japan interactions achieves conventional significance (p < 0.05) under wild bootstrap",
    "N = 4 country clusters means even the wild bootstrap has severely limited statistical power",
    "Interpret as: no robust causal evidence of a positive Japan governance re-rating effect at country-panel level",
    "The event-window and synthetic control corroborate structural descriptive trends; "
    "none of the three designs delivers clean causal identification",
]
add_para(tf2, "Interpretation", size=17, bold=True)
for t in takeaways:
    add_para(tf2, "•  " + t, size=14, space_before=3, indent=1)


# ─── Slide: Geopolitical Risk Results + Figure 3 ────────────────────────
sl = add_slide()
title_bar(sl, "Geopolitical Risk Sub-Analysis",
          subtitle="GPR escalation coefficient is statistically insignificant — risk priced as a level discount")

add_figure(sl,
           f"{FIGURES}/figure3_geo_risk.png",
           Inches(0.3), Inches(1.35),
           Inches(7.5), Inches(5.85))

rtf = txb(sl, Inches(8.0), Inches(1.35), Inches(5.0), Inches(5.85))
add_para(rtf, "Key results", size=17, bold=True)
geo_res = [
    "GPR escalation coefficient: –0.02 (t = –0.84, p = 0.40)",
    "Statistically insignificant at conventional levels",
    "TOPIX P/B coefficient: +0.63 (t = 5.40) — captures broader governance valuation convergence",
]
for g in geo_res:
    add_para(rtf, "•  " + g, size=14, space_before=5, indent=1)

add_para(rtf, "\nInterpretation", size=16, bold=True, space_before=10)
interps = [
    "Markets appear to price North Korean geopolitical risk continuously as a structural level discount, "
    "not through discrete monthly reactions to escalation events",
    "Consistent with IMF (2021): provocations primarily affect systematic risk / market beta",
    "Caveat: GPR dummy may be correlated with global risk-off episodes; monthly frequency may miss ephemeral reactions",
]
for i in interps:
    add_para(rtf, "•  " + i, size=13, space_before=4, indent=1)


# ─── Slide: Synthetic Control Results ───────────────────────────────────
sl = add_slide()
title_bar(sl, "Synthetic Control: 2023 TSE P/B Reform",
          subtitle="Actual minus synthetic Japan P/B; donor pool = regional Asian EM peers")

add_figure(sl,
           f"{FIGURES}/figure_synth_gap.png",
           Inches(0.3), Inches(1.35),
           Inches(7.5), Inches(5.85))

rtf = txb(sl, Inches(8.0), Inches(1.35), Inches(5.0), Inches(5.85))
add_para(rtf, "Construction", size=17, bold=True)
constr = [
    "Donor pool: MSCI EM, MSCI EM Asia, MSCI HK, MSCI Taiwan, MSCI China, MSCI India, MSCI Indonesia",
    "Optimization on demeaned TOPIX P/B (each unit demaned by pre-treatment mean)",
    "Demeaning resolves convex-hull boundary from Japan's structurally low absolute level",
    "KOSPI excluded to avoid using the primary comparison market as part of the Japan counterfactual",
]
for c in constr:
    add_para(rtf, "•  " + c, size=13, space_before=4, indent=1)

add_para(rtf, "\nResults", size=17, bold=True, space_before=8)
sc_res = [
    "Optimal weight: MSCI EM Asia = 1.0 (all other donors: 0%) — single-donor solution",
    "Pre-treatment RMSPE: 0.1451 (close to but below 0.15 threshold)",
    "Average post-reform gap: modestly positive (Japan slightly above synthetic)",
    "First 18-month gap change: –0.0028 per month (fading, not widening)",
]
for r in sc_res:
    add_para(rtf, "•  " + r, size=13, space_before=4, indent=1)

add_para(rtf,
    "\nConclusion: Japan's post-2023 P/B lift broadly matched its synthetic Asian EM peer. "
    "No sustained monotonic Japan-specific governance premium detected after 2023.",
    size=13, italic=True, color=MGRAY, space_before=8)


# ═══════════════════════════════════════════════════════════════════════════
#  SECTION: DISCUSSION
# ═══════════════════════════════════════════════════════════════════════════
section_divider("Discussion", "Limitations · Policy Implications · Conclusions")


# ─── Slide: Limitations ─────────────────────────────────────────────────
sl = add_slide()
title_bar(sl, "Key Limitations",
          subtitle="Honest assessment of what the evidence can and cannot support")

tf = body_tf(sl, top=Inches(1.35))

limits = [
    ("1.  Single Treated Unit  (N = 1 treated country)",
     [
         "Japan is the only treated country → cluster-robust SEs for the treated unit are undefined",
         "Three reform cohorts are not independent: all three events occurred to the same country",
         "2014 and 2015 reform windows overlap → cannot be treated as separate experiments",
         "Event-window evidence should be read as descriptive timing evidence, not precise causal estimates",
     ]),
    ("2.  The Abenomics Confound",
     [
         "Japan's reforms coincided with Abenomics: BoJ QQE (Apr 2013) + Yield Curve Control (Sep 2016)",
         "Yen depreciation → improved exporter competitiveness + compressed discount rates → mechanically higher P/B",
         "2023 TSE reform occurs outside peak Abenomics → cleanest governance-only window",
         "Two-way FE absorbs common time shocks but cannot cleanly separate governance reform from Japan-specific macro",
     ]),
    ("3.  Japan-to-Korea Generalizability",
     [
         "Korea's chaebol concentration more extreme than Japan's keiretsu structure",
         "If controlling families can dilute reform implementation, Korea's trajectory may be slower or less complete",
         "Geopolitical risk premium from North Korea has no Japanese equivalent — governance reform alone cannot eliminate it",
         "Counterfactual projection captures only the governance-reform component of the Korea Discount",
     ]),
]

for lim_title, bullets in limits:
    add_para(tf, lim_title, size=15, bold=True, space_before=6)
    for b in bullets:
        add_para(tf, "•  " + b, size=13, space_before=2, indent=1)


# ─── Slide: Policy Implications ─────────────────────────────────────────
sl = add_slide()
title_bar(sl, "Policy Implications",
          subtitle="Three sequenced levers — calibrated from Japan's reform experience (grounded in institutional channels, not precise causal estimates)")

tf = body_tf(sl, top=Inches(1.35))

recs = [
    ("1.  FSC Mandatory Capital Efficiency Disclosure",
     "Require all KOSPI companies trading below 1.0× book value to disclose multi-year P/B improvement plans "
     "on a mandatory comply-or-explain basis. Korea's 2024 Value-Up Program moves in this direction but "
     "remains voluntary — compliance with voluntary programs is historically lower. "
     "Disclosures should include: cost-of-equity baseline, specific P/B / ROE / dividend targets, "
     "operational action plans, and annual board-level progress reports with KRX-enforced penalties."),
    ("2.  KRX Listing Standards Reform",
     "Amend KRX rules: (i) ≥ 50% independent board composition for large-cap KOSPI issuers; "
     "(ii) enhanced disclosure of material related-party transactions (NYSE FPI thresholds); "
     "(iii) audit committee independence from controlling shareholders for chaebol-affiliated companies. "
     "Directly reduces principal–principal agency costs (La Porta et al. 1998, 2002)."),
    ("3.  Korean Stewardship Code Strengthening",
     "Amend 2016 Stewardship Code: (i) annual publication of voting records with rationale for votes against "
     "management; (ii) mandatory engagement with investee companies showing persistent below-peer P/B for "
     "≥ 3 years; (iii) NPS to establish a dedicated governance engagement division. "
     "Japan's cross-shareholding unwinding (Miyajima 2023) shows Stewardship Code activation "
     "is a key transmission mechanism between governance reform and P/B improvement."),
]

for title, body in recs:
    add_para(tf, title, size=15, bold=True, space_before=8)
    add_para(tf, body, size=13, space_before=2, indent=1)


# ─── Slide: Counterfactual Projection + Figure 4 ────────────────────────
sl = add_slide()
title_bar(sl, "Illustrative Counterfactual Projection",
          subtitle="KOSPI P/B under a Korea governance reform scenario — stress test, not a forecast")

add_figure(sl,
           f"{FIGURES}/figure4_counterfactual_projection.png",
           Inches(0.3), Inches(1.35),
           Inches(7.5), Inches(5.85))

rtf = txb(sl, Inches(8.0), Inches(1.35), Inches(5.0), Inches(5.85))
add_para(rtf, "Construction", size=16, bold=True)
proj = [
    "Starting point: KOSPI P/B ≈ 0.93× (December 2024)",
    "Applied rate: average monthly change in Japan's synthetic-control gap over first 18 post-reform months (–0.0028 per month)",
    "Uncertainty band: ± RMSPE = ± 0.1451",
    "Shaded region represents synthetic-control estimation uncertainty",
]
for p in proj:
    add_para(rtf, "•  " + p, size=13, space_before=4, indent=1)

add_para(rtf, "\nResult", size=16, bold=True, space_before=8)
add_para(rtf,
    "Because the Japan post-2023 relative gap change is negative (–0.0028/month), "
    "the mechanical projection does not close the Korea Discount. "
    "This is a stress test of the Japan-calibrated extrapolation, not a prediction of policy success.",
    size=13, space_before=2)

add_para(rtf, "\nCaveats", size=16, bold=True, space_before=8)
caveats = [
    "Assumes Korea replicates Japan's post-2023 P/B trajectory",
    "Cannot resolve the North Korea geopolitical risk premium",
    "FSC enforcement capacity and KRX compliance culture differ from Japan's",
    "Order-of-magnitude policy benchmark  ·  NOT a forecast",
]
for c in caveats:
    add_para(rtf, "•  " + c, size=12, space_before=3, italic=True, color=MGRAY, indent=1)


# ─── Slide: Conclusion ──────────────────────────────────────────────────
sl = add_slide()
title_bar(sl, "Conclusion")

tf = body_tf(sl, top=Inches(1.35))

sections = [
    ("What we document",
     [
         "KOSPI P/B averages –0.177× below TOPIX (t = –3.23) and –0.601× below MSCI EM (t = –10.30) over 2004–2024",
         "Discount has widened post-2023: KOSPI ≈ 0.93× vs. TOPIX ≈ 1.36×",
         "Three structural channels: chaebol opacity, minority-shareholder recourse deficit, geopolitical risk premium",
     ]),
    ("What we find empirically",
     [
         "Event windows show large spread movements around Japan's reform dates; not causally identified",
         "Panel OLS reform interactions: p-values 0.625, 0.375, 0.500 — noisy, not robust evidence of causal re-rating",
         "Synthetic control: RMSPE = 0.1451; full weight to MSCI EM Asia; post-reform gap fades rather than widens",
         "Geopolitical risk coefficient: –0.02 (t = –0.84, p = 0.40) — priced as structural level discount",
     ]),
    ("Directions for future research",
     [
         "Firm-level chaebol governance panel to complement country-level time-series identification",
         "Event study of Korea's Value-Up Program as multi-year post-announcement data accumulate (2024–2026+)",
         "Cross-country extension: China A-shares, India, Gulf — do Korea's structural channels generalize?",
     ]),
]

for sec_title, bullets in sections:
    add_para(tf, sec_title, size=16, bold=True, space_before=8)
    for b in bullets:
        add_para(tf, "•  " + b, size=13, space_before=2, indent=1)

add_para(tf,
    "\nBottom line: The Korea Discount is persistent and statistically significant. "
    "Japan's governance reform experience provides useful policy benchmarks but does not, "
    "by itself, identify the valuation gain Korea would realize from adopting Japan-style reforms.",
    size=14, italic=True, color=MGRAY, space_before=8)


# ═══════════════════════════════════════════════════════════════════════════
#  Save
# ═══════════════════════════════════════════════════════════════════════════
out = "/Users/dandan/Desktop/Projects/kor-discount/output/korea_discount_slides.pptx"
prs.save(out)
print(f"Saved → {out}")
